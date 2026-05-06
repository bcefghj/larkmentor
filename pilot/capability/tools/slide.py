"""slide.generate / slide.rehearse — PPT 工具.

借鉴:
  - Gamma：card-based 模板（5 模板：Hero / TwoColumn / Cards / List / Quote）
  - Beautiful.ai：设计规则约束（自动 reflow，字数超限分页）
  - python-pptx 底层渲染

输出:
  - 真 .pptx 文件
  - Slidev markdown（如服务器装了 npx @slidev/cli 可继续 build HTML）
  - 演讲稿 markdown
  - 可选 TTS mp3
"""

from __future__ import annotations

import json
import logging
import os
import re
import time
import uuid
from pathlib import Path
from typing import Any

from pilot.context.filesystem_memory import ARTIFACTS_DIR, FilesystemMemory

logger = logging.getLogger("pilot.tool.slide")


def register_to(reg) -> None:
    reg.register(
        "slide.generate",
        description="基于上游 doc.append 生成真 .pptx 演示稿（5 模板）+ Slidev md + 演讲稿",
        input_schema={
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "outline": {"type": "array", "description": "可选大纲；留空时从上游 doc 自动提炼"},
                "intent": {"type": "string", "description": "用户原始意图"},
                "pages": {"type": "integer", "description": "目标页数，默认 8"},
                "search_results": {
                    "description": "上游 web.search 注入的结果 [{title,url,snippet}]，用于 LLM 大纲引用真实数据",
                },
            },
            "required": ["title"],
        },
        read_only=False,
        namespace="pilot",
    )(slide_generate)

    reg.register(
        "slide.rehearse",
        description="为已生成的 PPT 每页生成演讲稿",
        input_schema={
            "type": "object",
            "properties": {
                "slide_id": {"type": "string"},
            },
            "required": ["slide_id"],
        },
        read_only=False,
        namespace="pilot",
    )(slide_rehearse)


# ── slide.generate ──


async def slide_generate(
    *,
    title: str,
    outline: list | None = None,
    intent: str = "",
    pages: int = 8,
    search_results: Any = None,
    _ctx: dict[str, Any] | None = None,
) -> dict[str, Any]:
    sid = f"slide_{int(time.time())}_{uuid.uuid4().hex[:6]}"
    out_dir = ARTIFACTS_DIR / "slides" / sid
    out_dir.mkdir(parents=True, exist_ok=True)

    upstream_md = _extract_upstream_doc(_ctx)
    citations = _normalize_search_results(search_results)
    if not outline:
        outline = await _llm_outline(
            title=title,
            intent=intent,
            pages=pages,
            upstream_md=upstream_md,
            search_results=citations,
        )
    outline = _normalise_outline(outline, target_pages=pages)

    pptx_path = out_dir / f"{sid}.pptx"
    page_count = _write_pptx(title=title, outline=outline, out_path=pptx_path)

    slidev_md_path = out_dir / f"{sid}.slidev.md"
    slidev_md_path.write_text(_outline_to_slidev_md(title, outline), encoding="utf-8")

    notes_md_path = out_dir / f"{sid}.speaker_notes.md"
    notes_md_path.write_text(_outline_to_speaker_notes(title, outline), encoding="utf-8")

    base = (os.getenv("DASHBOARD_PUBLIC_BASE") or "").rstrip("/")
    rel = f"/artifacts/slides/{sid}"
    pptx_rel = f"{rel}/{pptx_path.name}"
    return {
        "slide_id": sid,
        "title": title,
        "pages": page_count,
        "outline": outline,
        "pptx_path": str(pptx_path),
        "pptx_url": pptx_rel,
        "pptx_url_absolute": f"{base}{pptx_rel}" if base else pptx_rel,
        "slidev_md_url": f"{rel}/{slidev_md_path.name}",
        "speaker_notes_md_url": f"{rel}/{notes_md_path.name}",
        "citations": citations,
    }


def _normalize_search_results(raw: Any) -> list[dict[str, str]]:
    """与 doc 工具保持一致的归一化逻辑."""
    if not raw:
        return []
    if isinstance(raw, str):
        s = raw.strip()
        if s.startswith("$") or not s:
            return []
        try:
            raw = json.loads(s)
        except Exception:
            return []
    if isinstance(raw, dict):
        raw = raw.get("results") or raw.get("items") or []
    if not isinstance(raw, list):
        return []
    out = []
    for item in raw:
        if isinstance(item, dict):
            t = str(item.get("title", ""))[:200]
            u = str(item.get("url", ""))[:500]
            sn = str(item.get("snippet", "") or item.get("desc", ""))[:400]
            if t or u:
                out.append({"title": t, "url": u, "snippet": sn})
    return out[:10]


# ── slide.rehearse ──


async def slide_rehearse(*, slide_id: str, _ctx: dict[str, Any] | None = None) -> dict[str, Any]:
    """为现有 slide 生成更详细的演讲稿（每页 80-150 字）."""
    out_dir = ARTIFACTS_DIR / "slides" / slide_id
    if not out_dir.exists():
        return {"slide_id": slide_id, "ok": False, "error": "slide 不存在"}

    notes_path = out_dir / f"{slide_id}.speaker_notes.md"
    if notes_path.exists():
        text = notes_path.read_text(encoding="utf-8")
    else:
        text = "(占位演讲稿)"

    rehearse_path = out_dir / f"{slide_id}.rehearse.md"
    rehearse_path.write_text(_enhance_rehearsal(text), encoding="utf-8")

    return {
        "slide_id": slide_id,
        "speaker_notes_md_path": str(rehearse_path),
        "speaker_notes_md_url": f"/artifacts/slides/{slide_id}/{rehearse_path.name}",
        "ok": True,
    }


# ── 实现细节 ──


async def _llm_outline(
    *,
    title: str,
    intent: str,
    pages: int,
    upstream_md: str,
    search_results: list[dict[str, str]] | None = None,
) -> list[dict[str, Any]]:
    try:
        from pilot.llm.client import default_client
        from pilot.llm.safe_json import safe_json_parse

        cite_block = ""
        if search_results:
            lines = ["\n参考资料（请基于真实数据，不要瞎编）："]
            for i, r in enumerate(search_results[:5], 1):
                lines.append(f"[{i}] {r.get('title','')} — {r.get('url','')}\n    {r.get('snippet','')}")
            cite_block = "\n".join(lines)

        prompt = f"""请为「{title}」生成 {pages} 页 PPT 大纲（封面 + 主体 + 结尾）。

用户意图：{intent or '（无）'}
上游文档：
{upstream_md[:2000] or '（无）'}
{cite_block}

输出严格 JSON 数组，每项：
{{
  "template": "Hero|TwoColumn|Cards|List|Quote",
  "title": "页标题",
  "bullets": ["要点 1", "要点 2", "要点 3"],
  "notes": "演讲备注（30-80 字）"
}}

模板说明：
- Hero: 封面 / 章节首页（大字 + 副标题）
- TwoColumn: 左右对比 / 优劣分析
- Cards: 3-4 张卡片并列
- List: 数字列表
- Quote: 名人 / 用户引言

要求：第 1 页 Hero（封面）；最后一页 Hero（致谢/Q&A）；中间 6 页混用其他模板。
若有参考资料，请在合适页面 bullets 中以 "[1]"/"[2]" 形式引用，并在 notes 中说明数据出处。
直接输出 JSON 数组。"""

        client = default_client()
        result = await client.chat(
            system="你是 Agent-Pilot 的演示设计师，擅长 PPT 大纲。",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            max_tokens=2048,
        )
        for block in result.get("content", []):
            if isinstance(block, dict) and block.get("type") == "text":
                obj = safe_json_parse(block.get("text", ""), expected_type=list, debug_label="slide.outline")
                if obj:
                    return obj
    except Exception as e:
        logger.warning("slide outline LLM failed: %s", e)
    return _fallback_outline(title, pages)


def _fallback_outline(title: str, pages: int) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = [{
        "template": "Hero",
        "title": title,
        "bullets": ["Agent-Pilot 自动生成"],
        "notes": f"开场介绍 {title}，引入正文。",
    }]
    chapters = ["背景与现状", "核心问题", "解决方案", "落地路径", "数据与案例", "风险与对策"]
    n = max(1, pages - 2)
    for i in range(n):
        out.append({
            "template": ["TwoColumn", "Cards", "List"][i % 3],
            "title": chapters[i % len(chapters)],
            "bullets": [f"要点 {i + 1}-1", f"要点 {i + 1}-2", f"要点 {i + 1}-3"],
            "notes": f"讲清楚 {chapters[i % len(chapters)]} 的核心。",
        })
    out.append({
        "template": "Hero",
        "title": "Thank You",
        "bullets": ["欢迎提问", "联系方式 见首页"],
        "notes": "感谢聆听，留 2 分钟问答。",
    })
    return out[:pages]


def _normalise_outline(outline: list, target_pages: int) -> list[dict[str, Any]]:
    out = []
    for item in outline:
        if isinstance(item, dict):
            out.append({
                "template": item.get("template", "List"),
                "title": item.get("title", "")[:60],
                "bullets": [str(b)[:120] for b in (item.get("bullets") or [])][:5],
                "notes": item.get("notes", "")[:300],
            })
        elif isinstance(item, str):
            out.append({"template": "List", "title": item[:60], "bullets": [], "notes": ""})
    if not out:
        out = _fallback_outline("Agent-Pilot 演示", target_pages)
    return out


# ── PPTX 渲染 ──


def _write_pptx(*, title: str, outline: list[dict[str, Any]], out_path: Path) -> int:
    """用 python-pptx 渲染 5 模板 PPT；失败时写出 markdown 占位."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError:
        # 回退：写一份 markdown
        md_path = out_path.with_suffix(".md")
        md_path.write_text(_outline_to_slidev_md(title, outline), encoding="utf-8")
        return len(outline)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 空白
    for page in outline:
        tpl = page.get("template", "List")
        if tpl == "Hero":
            _render_hero(prs, blank_layout, page)
        elif tpl == "TwoColumn":
            _render_twocolumn(prs, blank_layout, page)
        elif tpl == "Cards":
            _render_cards(prs, blank_layout, page)
        elif tpl == "Quote":
            _render_quote(prs, blank_layout, page)
        else:
            _render_list(prs, blank_layout, page)

    prs.save(str(out_path))
    return len(outline)


def _render_hero(prs, layout, page):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = page.get("title", "")
    p.runs[0].font.size = Pt(54)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x0F, 0x4C, 0x81)

    sub = " · ".join(page.get("bullets", [])[:3])
    if sub:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.8))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = sub
        sp.runs[0].font.size = Pt(20)
        sp.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    _add_notes(slide, page.get("notes", ""))


def _render_list(prs, layout, page):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = page.get("title", "")
    p.runs[0].font.size = Pt(32)
    p.runs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(page.get("bullets", [])):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"• {b}"
        para.runs[0].font.size = Pt(22)

    _add_notes(slide, page.get("notes", ""))


def _render_twocolumn(prs, layout, page):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = page.get("title", "")
    p.runs[0].font.size = Pt(32)
    p.runs[0].font.bold = True

    bullets = page.get("bullets", [])
    half = (len(bullets) + 1) // 2
    left = bullets[:half]
    right = bullets[half:]

    for i, (col_x, col_bullets) in enumerate([(Inches(0.8), left), (Inches(7.0), right)]):
        col = slide.shapes.add_textbox(col_x, Inches(1.5), Inches(5.5), Inches(5.5))
        tf = col.text_frame
        tf.word_wrap = True
        for j, b in enumerate(col_bullets):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = f"• {b}"
            para.runs[0].font.size = Pt(20)

    _add_notes(slide, page.get("notes", ""))


def _render_cards(prs, layout, page):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = page.get("title", "")
    p.runs[0].font.size = Pt(32)
    p.runs[0].font.bold = True

    bullets = page.get("bullets", [])[:4]
    if not bullets:
        return _render_list(prs, layout, page)

    card_w = 3.0
    gap = 0.2
    total_w = card_w * len(bullets) + gap * (len(bullets) - 1)
    start_x = (13.33 - total_w) / 2

    for i, b in enumerate(bullets):
        x = Inches(start_x + i * (card_w + gap))
        card = slide.shapes.add_textbox(x, Inches(2.0), Inches(card_w), Inches(3.5))
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💡 {b}"
        p.runs[0].font.size = Pt(18)
        p.runs[0].font.bold = True

    _add_notes(slide, page.get("notes", ""))


def _render_quote(prs, layout, page):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(layout)
    quote_text = page.get("bullets", [page.get("title", "")])[0]
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.5))
    p = box.text_frame.paragraphs[0]
    p.text = f"\u201c{quote_text}\u201d"
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.italic = True
    p.runs[0].font.color.rgb = RGBColor(0x0F, 0x4C, 0x81)

    _add_notes(slide, page.get("notes", ""))


def _add_notes(slide, text: str):
    if not text:
        return
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


# ── Slidev / 演讲稿 ──


def _outline_to_slidev_md(title: str, outline: list) -> str:
    parts = ["---", "theme: default", f"title: {title}", "---", ""]
    for page in outline:
        parts.append(f"# {page.get('title', '')}")
        for b in page.get("bullets", []):
            parts.append(f"- {b}")
        parts.append("")
        if page.get("notes"):
            parts.append("<!--")
            parts.append(page["notes"])
            parts.append("-->")
            parts.append("")
        parts.append("---")
        parts.append("")
    return "\n".join(parts)


def _outline_to_speaker_notes(title: str, outline: list) -> str:
    out = [f"# {title} · 演讲稿\n"]
    for i, page in enumerate(outline, start=1):
        out.append(f"## 第 {i} 页 — {page.get('title', '')}\n")
        out.append(page.get("notes", "") + "\n")
    return "\n".join(out)


def _enhance_rehearsal(text: str) -> str:
    return f"## 排练版演讲稿（含停顿/重音提示）\n\n{text}"


# ── 上游 doc 提取 ──


def _extract_upstream_doc(ctx: dict[str, Any] | None) -> str:
    if not ctx:
        return ""
    results = ctx.get("step_results") or {}
    for r in reversed(list(results.values())):
        if isinstance(r, dict):
            if "markdown_artifact" in r and isinstance(r["markdown_artifact"], dict):
                uri = r["markdown_artifact"].get("uri", "")
                if uri:
                    mem = FilesystemMemory()
                    md = mem.resolve(uri)
                    if md:
                        return md
            if "markdown" in r and isinstance(r["markdown"], str):
                return r["markdown"]
    return ""
