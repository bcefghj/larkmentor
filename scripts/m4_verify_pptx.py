#!/usr/bin/env python3
"""M4 verification: generate a real PPTX with mock outline (no LLM) and inspect it.

Usage:
    python3 scripts/m4_verify_pptx.py
"""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from agent_pilot.tools.slide import slide_generate, _make_pptx, ARTIFACTS_DIR


def main():
    # Use a hand-crafted outline so we don't depend on LLM
    title = "AI Agent 发展趋势报告"
    outline = [
        {"title": title, "bullets": ["AI 驱动办公协同的下一站", "Agent-Pilot 团队"], "note": "开场介绍"},
        {"title": "目录", "bullets": ["产业现状", "技术演进", "应用场景", "风险与展望"], "note": "今天讲四部分"},
        {"title": "产业现状", "bullets": [
            "全球 AI Agent 市场 2024 年突破 50 亿美元",
            "头部模型公司加速 Agent 化布局（OpenAI / Anthropic / 智谱）",
            "国内厂商飞书/钉钉/企微 IM 内开始集成 AI Agent",
        ], "note": "市场规模与玩家"},
        {"title": "技术演进", "bullets": [
            "从单轮 chat 到 ReAct/Plan-Execute 多步规划",
            "LLM 函数调用 + MCP 协议构建工具生态",
            "Multi-Agent 协作（Builder-Validator）成为新标准",
        ], "note": "技术路线"},
        {"title": "应用场景", "bullets": [
            "办公协同：从 IM 对话到文档/PPT 自动闭环",
            "代码助手：Cursor / Copilot 重塑研发工作流",
            "客户服务：7×24 多模态智能客服",
        ], "note": "三大落地场景"},
        {"title": "风险与展望", "bullets": [
            "幻觉与错误传播：需要 Critic Agent 二次校验",
            "数据隐私：本地化部署 + 端到端加密",
            "未来：AI Native 应用全面替代 GUI 堆砌",
        ], "note": "需要警惕"},
        {"title": "Thank You", "bullets": [
            "Agent-Pilot · 让对话直接变成产物",
            "GitHub: bcefghj/Agent-Pilot",
            "戴尚好 / 李洁盈",
        ], "note": "感谢评委"},
    ]

    out_path = ARTIFACTS_DIR / "m4_verify" / "demo.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    pages = _make_pptx(title, outline, out_path)

    print(f"✓ wrote {pages} slides to {out_path}")
    size_kb = out_path.stat().st_size / 1024
    print(f"✓ file size: {size_kb:.1f} KB")

    # Verify with python-pptx
    from pptx import Presentation
    prs = Presentation(str(out_path))
    actual_slides = list(prs.slides)
    print(f"✓ pptx parse: {len(actual_slides)} slides")
    for i, slide in enumerate(actual_slides, 1):
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
        title_text = (text_parts[0] if text_parts else "(empty)")[:40]
        notes_text = ""
        try:
            notes_text = slide.notes_slide.notes_text_frame.text[:60]
        except Exception:
            pass
        print(f"  Slide {i}: {title_text!r} | notes: {notes_text!r}")

    assert len(actual_slides) >= 6, f"Expected at least 6 slides, got {len(actual_slides)}"
    assert size_kb >= 8, f"Expected at least 8KB, got {size_kb} KB"
    print(f"\n✅ M4 PPTX verification PASSED")
    print(f"   Open with: open '{out_path}'")
    return out_path


if __name__ == "__main__":
    main()
