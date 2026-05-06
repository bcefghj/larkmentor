"""slide.generate + slide.rehearse – v13 真 PPTX 三件套.

This is the single source of truth for "演示稿/PPT" generation. It produces:

1. **`.pptx`** via ``python-pptx`` – a real PowerPoint file with cover, TOC,
   content slides (title + bullets + speaker notes), and a Thank-You page.
2. **Slidev HTML** – a Slidev-formatted markdown that can be built into a
   static slideshow with ``npx @slidev/cli build``. Falls back to a
   Markdown-only artifact when ``npx`` isn't available.
3. **Speaker-notes Markdown** – a per-page narration script. (The Markdown is
   always produced; TTS mp3 is best-effort.)

Optional: TTS mp3 (gTTS) when ``AGENT_PILOT_ENABLE_TTS=1`` is set and the
network can reach Google.

The tool also creates a **preview Feishu Docx** (one heading + bullets per
slide) so users on iOS/Android can flip through it natively. The .pptx file
itself is uploaded to Feishu cloud Drive when an ``FEISHU_FOLDER_TOKEN`` is
configured; otherwise it's served via the Dashboard ``/artifacts/slides/...``
static route.
"""

from __future__ import annotations

import json
import logging
import os
import shutil
import subprocess
import time
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from agent_pilot.llm.safe_json import safe_json_parse

logger = logging.getLogger("agent_pilot.tool.slide")

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
ARTIFACTS_DIR = PROJECT_ROOT / "data" / "artifacts" / "slides"


def _ensure_dir(p: Path) -> None:
    p.mkdir(parents=True, exist_ok=True)


# ── Public entry: slide.generate ──────────────────────────────────────────────


def slide_generate(step, ctx: Dict[str, Any]) -> Dict[str, Any]:
    """Generate the PPT trio for the given context.

    Returns a dict with:
      - slide_id: short identifier
      - title: the slide deck title
      - pages: number of content slides
      - outline: the structured outline (used by slide.rehearse)
      - pptx_path / pptx_url: real .pptx file
      - slidev_md_path / slidev_html_path / slidev_url: Slidev HTML preview
      - speaker_notes_md_path / speaker_notes_md_url: per-page script
      - tts_mp3_path / tts_url: optional voiceover
      - preview_doc_url / preview_doc_token / feishu_url: Feishu Docx preview
      - source: "feishu+pptx" | "pptx_only"
    """
    args = ctx.get("resolved_args") or {}
    title = (args.get("title") or "").strip()
    if not title or "{{" in title:
        title = _title_from_ctx(ctx)
    outline = args.get("outline")
    if not outline or (isinstance(outline, str) and "{{" in outline):
        outline = _build_outline_from_doc_or_intent(ctx, title)
    outline = _normalise_outline(outline)
    if len(outline) < 4:  # pad to a usable minimum
        outline = _pad_outline(outline, title)

    slide_id = f"slide_{int(time.time())}_{uuid.uuid4().hex[:6]}"
    slide_dir = ARTIFACTS_DIR / slide_id
    _ensure_dir(slide_dir)

    # 1. Real .pptx
    pptx_path = slide_dir / f"{slide_id}.pptx"
    pages = _make_pptx(title, outline, pptx_path)
    logger.info("slide.generate: wrote pptx=%s pages=%d", pptx_path, pages)

    # 2. Slidev HTML (best-effort; markdown always emitted)
    slidev_md_path = slide_dir / f"{slide_id}.slidev.md"
    slidev_md_path.write_text(_outline_to_slidev_md(title, outline), encoding="utf-8")
    slidev_html_path = _build_slidev_html(slidev_md_path, slide_dir)

    # 3. Speaker notes markdown
    notes_md_path = slide_dir / f"{slide_id}.speaker_notes.md"
    notes_md_path.write_text(_outline_to_speaker_notes(title, outline), encoding="utf-8")

    # 4. Optional TTS mp3 (only when explicitly enabled)
    tts_mp3_path = _make_tts_optional(title, outline, slide_dir, slide_id)

    # 5. Feishu Docx preview (always good to have for mobile users)
    feishu_preview = _create_feishu_preview_doc(title, outline)

    # 6. Try uploading the .pptx to Feishu Drive
    pptx_feishu_url = _upload_pptx_to_feishu_drive(pptx_path)

    base_artifact_url = f"/artifacts/slides/{slide_id}"
    result: Dict[str, Any] = {
        "slide_id": slide_id,
        "title": title,
        "pages": pages,
        "outline": outline,
        "pptx_path": str(pptx_path),
        "pptx_url": pptx_feishu_url or f"{base_artifact_url}/{pptx_path.name}",
        "slidev_md_path": str(slidev_md_path),
        "speaker_notes_md_path": str(notes_md_path),
    }
    if slidev_html_path:
        result["slidev_html_path"] = str(slidev_html_path)
        result["slidev_url"] = f"{base_artifact_url}/index.html"
    if tts_mp3_path:
        result["tts_mp3_path"] = str(tts_mp3_path)
        result["tts_url"] = f"{base_artifact_url}/{tts_mp3_path.name}"

    if feishu_preview.get("url"):
        result["preview_doc_url"] = feishu_preview["url"]
        result["preview_doc_token"] = feishu_preview.get("doc_token", "")
        result["feishu_url"] = feishu_preview["url"]
        result["source"] = "feishu+pptx" if pptx_feishu_url else "feishu+pptx_local"
    else:
        result["source"] = "pptx_only"

    return result


def slide_rehearse(step, ctx: Dict[str, Any]) -> Dict[str, Any]:
    """Generate professional speaker notes via a dedicated LLM call."""
    args = ctx.get("resolved_args") or {}
    slide_id = args.get("slide_id") or ""
    outline = args.get("outline") or []

    if (not outline) and slide_id:
        d = ARTIFACTS_DIR / slide_id
        for child in [d / f"{slide_id}.slidev.md"]:
            if child.exists():
                outline = _parse_slidev_md(child)
                break

    notes = _llm_rehearsal_notes(outline)
    if not notes:
        notes = [
            {
                "page": i,
                "title": page.get("title", ""),
                "speaker_note": _speaker_note_for_page(page),
                "duration_sec": 45,
                "transition": "",
            }
            for i, page in enumerate(outline or [], start=1)
        ]

    if slide_id:
        notes_path = ARTIFACTS_DIR / slide_id / f"{slide_id}.rehearsal.md"
        _ensure_dir(notes_path.parent)
        lines = ["# 演讲排练稿\n"]
        for n in notes:
            lines.append(f"## 第 {n['page']} 页：{n['title']}  (建议 {n['duration_sec']}s)\n")
            lines.append(n["speaker_note"])
            if n.get("transition"):
                lines.append(f"\n> 过渡：{n['transition']}\n")
            lines.append("")
        notes_path.write_text("\n".join(lines), encoding="utf-8")

    return {
        "slide_id": slide_id,
        "rehearsal_pages": len(notes),
        "total_duration_sec": sum(n.get("duration_sec", 45) for n in notes),
        "speaker_notes": notes,
    }


def _llm_rehearsal_notes(outline: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not outline:
        return []
    try:
        from llm.llm_client import chat, LLM_FALLBACK_MSG
    except ImportError:
        return []

    pages_summary = "\n".join(
        f"第{i}页「{p.get('title','')}」要点：{'；'.join((p.get('bullets') or [])[:3])}"
        for i, p in enumerate(outline, 1)
    )
    prompt = f"""你是一位演讲教练。请为以下 PPT 的每一页编写自然口语化的演讲稿。

## PPT 大纲
{pages_summary}

请输出 JSON 数组，每个元素：
{{"page": 1, "title": "页面标题", "speaker_note": "120-180字的自然口语演讲稿，像真人在台上讲", "duration_sec": 45, "transition": "到下一页的过渡语（15字以内）"}}

要求：
1. 演讲稿 120-180 字，自然口语，不要书面语
2. 封面页简短开场（60字），Thank You 页简短致谢（60字）
3. 内容页要有逻辑过渡，引用要点中的数据
4. duration_sec 根据内容长度给出合理估计（30-90秒）
5. transition 是过渡到下一页的一句话
直接输出 JSON 数组。"""

    raw = chat(prompt, temperature=0.5, max_tokens=6000)
    if not raw or raw.strip() == LLM_FALLBACK_MSG:
        return []
    arr = safe_json_parse(raw, expected_type=list, debug_label="rehearsal")
    if not arr:
        return []
    result = []
    for item in arr:
        if not isinstance(item, dict):
            continue
        result.append({
            "page": int(item.get("page", len(result) + 1)),
            "title": str(item.get("title", "")),
            "speaker_note": str(item.get("speaker_note", "")),
            "duration_sec": int(item.get("duration_sec", 45)),
            "transition": str(item.get("transition", "")),
        })
    return result


# ── Outline derivation ────────────────────────────────────────────────────────


def _title_from_ctx(ctx: Dict[str, Any]) -> str:
    intent = (
        ctx.get("original_intent")
        or ctx.get("intent")
        or "Agent-Pilot 演示"
    )
    intent = intent.replace("\n", " ").strip()
    return intent[:30] or "Agent-Pilot 演示"


def _build_outline_from_doc_or_intent(ctx: Dict[str, Any], title: str) -> List[Dict[str, Any]]:
    """LLM with doc → LLM with title → template."""
    doc_md = ""
    step_results = ctx.get("step_results") or {}
    for r in step_results.values():
        if isinstance(r, dict) and r.get("markdown_content"):
            doc_md = r["markdown_content"]
            break
    intent = ctx.get("original_intent") or ctx.get("intent") or title

    outline = _llm_outline(title, intent, doc_md)
    if outline and len(outline) >= 4:
        return outline

    logger.info("slide.outline: LLM gave %d pages, padding from template", len(outline))
    return _pad_outline(outline, title)


def _llm_outline(title: str, intent: str, doc_markdown: str) -> List[Dict[str, Any]]:
    try:
        from llm.llm_client import chat, LLM_FALLBACK_MSG
    except ImportError:
        return []

    if doc_markdown:
        preview = doc_markdown[:14000]
        prompt = f"""你是一位专业的演示稿设计专家。基于以下文档内容，提炼出一份 8-12 页的高质量 PPT 大纲。

## 用户需求
{intent}

## 已有文档内容（请基于此提炼 PPT，与文档保持一致）
{preview}

请以 JSON 数组格式输出，每个元素包含 title、bullets、note 三个字段：
[
  {{"title": "封面", "bullets": ["副标题（一句话价值主张）", "演讲者：Agent-Pilot 团队"], "note": "开场白"}},
  {{"title": "目录", "bullets": ["章节1", "章节2", "章节3"], "note": "今天要讲的内容"}},
  ...
]

要求：
1. 8-12 页，结构完整，逻辑清晰
2. 每页 3-5 个 bullets，每条是一句完整有信息量的话（不是短语）
3. note 字段是这页的演讲稿（80-120 字），用于讲稿
4. 包含：封面、目录、背景与痛点、核心方案/分析（3-5 页深度展开）、案例或数据、风险与对策、总结展望、Thank You
5. PPT 内容必须与文档保持一致，是文档关键论点的提炼
6. 直接输出 JSON 数组，不要任何前缀或代码块包裹"""
    else:
        prompt = f"""你是一位专业的演示稿设计专家。请为以下主题设计一份 8-12 页的高质量 PPT 大纲。

主题：{title}
用户需求：{intent}

请以 JSON 数组格式输出，每个元素包含 title、bullets、note 三个字段：
[
  {{"title": "封面", "bullets": ["一句话价值主张", "Agent-Pilot 团队"], "note": "开场白"}},
  {{"title": "目录", "bullets": ["章节1", "章节2"], "note": "今天要讲的内容"}},
  ...
]

要求：8-12 页；每页 3-5 个完整句子的要点；note 是演讲稿（80-120 字）；
包含：封面、目录、背景与痛点、核心内容（3-5 页）、案例/数据、风险与对策、总结、Thank You。
直接输出 JSON 数组，不要任何前缀。"""

    raw = chat(prompt, temperature=0.5, max_tokens=8192)
    if not raw or raw.strip() == LLM_FALLBACK_MSG:
        return []
    arr = safe_json_parse(raw, expected_type=list, debug_label="slide.outline")
    if not arr:
        return []
    out: List[Dict[str, Any]] = []
    for i, p in enumerate(arr, 1):
        if not isinstance(p, dict):
            continue
        bullets_raw = p.get("bullets") or p.get("points") or []
        if isinstance(bullets_raw, str):
            bullets_raw = [b.strip() for b in bullets_raw.split("\n") if b.strip()]
        out.append({
            "title": str(p.get("title") or f"Slide {i}").strip(),
            "bullets": [str(b).strip() for b in bullets_raw if str(b).strip()],
            "note": str(p.get("note") or p.get("speaker_note") or "").strip(),
        })
    return out


def _pad_outline(outline: List[Dict[str, Any]], title: str) -> List[Dict[str, Any]]:
    template = [
        {"title": title, "bullets": ["AI 驱动的办公协同新范式", "Agent-Pilot 团队"], "note": "开场介绍我们的项目"},
        {"title": "目录", "bullets": ["背景与痛点", "解决方案", "技术架构", "实践案例", "未来展望"], "note": "今天要讲五个部分"},
        {"title": "背景与痛点", "bullets": [
            "知识工作者每天被 IM 打断 30+ 次，信息散落在聊天与文档中",
            "需求从对话产生，但要手动整理为方案、PPT、归档",
            "团队跨端协作时缺乏一致的进度可视化",
        ], "note": "团队协作的真实困境"},
        {"title": "解决方案", "bullets": [
            "Agent-Pilot 在飞书 IM 里主动识别任务",
            "自动驱动文档/画布/演示稿三种产物联动生成",
            "WebSocket 多端同步，移动+桌面+Web 实时一致",
        ], "note": "我们的产品定位"},
        {"title": "技术架构", "bullets": [
            "三闸门主动识别（规则+LLM+最小信息）",
            "DAG 编排 + 4-Agent 协作工坊",
            "真 .pptx / Mermaid+tldraw / Slidev 三件套",
        ], "note": "核心技术栈"},
        {"title": "实践案例", "bullets": [
            "案例一：群聊讨论 → 一句话生成方案+PPT",
            "案例二：手机扫码确认 → 桌面端编辑 PPT 实时同步",
            "案例三：模糊意图 Agent 主动澄清后再执行",
        ], "note": "真实使用场景"},
        {"title": "未来展望", "bullets": [
            "更深 AI Native：自主任务规划与持续学习",
            "更强生态：飞书白板/多维表格深度集成",
            "更广平台：iOS/Android/macOS/Windows 全端",
        ], "note": "下一步计划"},
        {"title": "Thank You", "bullets": [
            "Agent-Pilot · 让对话直接变成产物",
            "GitHub: https://github.com/bcefghj/Agent-Pilot",
            "戴尚好 / 李洁盈 · 飞书 AI 校园挑战赛",
        ], "note": "感谢评委倾听"},
    ]
    if outline:
        merged = {p["title"]: p for p in outline}
        for default in template:
            merged.setdefault(default["title"], default)
        return list(merged.values())
    return template


def _normalise_outline(outline: Any) -> List[Dict[str, Any]]:
    if not outline:
        return []
    if isinstance(outline, str):
        outline = [p.strip() for p in outline.split("\n") if p.strip()]
    out: List[Dict[str, Any]] = []
    for i, p in enumerate(outline, 1):
        if isinstance(p, dict):
            bullets_raw = p.get("bullets") or p.get("points") or []
            if isinstance(bullets_raw, str):
                bullets_raw = [b.strip() for b in bullets_raw.split("\n") if b.strip()]
            out.append({
                "title": str(p.get("title") or f"Slide {i}"),
                "bullets": [str(b) for b in bullets_raw],
                "note": str(p.get("note") or p.get("speaker_note") or ""),
            })
        elif isinstance(p, str):
            out.append({"title": p, "bullets": [], "note": ""})
    return out


# ── Real PPTX rendering ──────────────────────────────────────────────────────


_NAVY = (0x0A, 0x1F, 0x3B)
_NAVY_MID = (0x14, 0x3A, 0x6B)
_GOLD = (0xF0, 0xA5, 0x30)
_GOLD_LIGHT = (0xF7, 0xCC, 0x7F)
_LIGHT = (0xF5, 0xF5, 0xF5)
_WHITE = (0xFF, 0xFF, 0xFF)
_DARK = (0x2C, 0x2C, 0x2C)
_MUTED = (0x88, 0x88, 0x88)
_ACCENT_BLUE = (0x3B, 0x82, 0xF6)
_ACCENT_GREEN = (0x10, 0xB9, 0x81)
_ACCENT_PURPLE = (0x8B, 0x5C, 0xF6)
_ACCENT_RED = (0xEF, 0x44, 0x44)

_BULLET_ICONS = ["▸", "◆", "●", "■", "★"]
_SECTION_COLORS = [_ACCENT_BLUE, _ACCENT_GREEN, _ACCENT_PURPLE, _GOLD, _ACCENT_RED]


def _make_pptx(title: str, outline: List[Dict[str, Any]], out_path: Path) -> int:
    """Render a professional .pptx with varied layouts per slide type."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def _rect(slide, l, t, w, h, rgb, *, alpha=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        s.line.fill.background()
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(*rgb)
        return s

    def _text(slide, l, t, w, h, text, *, sz=18, bold=False, rgb=_DARK, align="left", font="Microsoft YaHei"):
        tx = slide.shapes.add_textbox(l, t, w, h)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*rgb)
        return tx

    def _numbered_bullets(slide, l, t, w, h, bullets, *, accent=_ACCENT_BLUE, sz=18):
        tx = slide.shapes.add_textbox(l, t, w, h)
        tf = tx.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(10)
            p.space_after = Pt(6)
            p.level = 0
            num = p.add_run()
            num.text = f"{i + 1}  "
            num.font.name = "Microsoft YaHei"
            num.font.size = Pt(sz + 4)
            num.font.bold = True
            num.font.color.rgb = RGBColor(*accent)
            body = p.add_run()
            body.text = b
            body.font.name = "Microsoft YaHei"
            body.font.size = Pt(sz)
            body.font.color.rgb = RGBColor(*_DARK)

    def _icon_bullets(slide, l, t, w, h, bullets, *, icon_color=_GOLD, sz=18):
        tx = slide.shapes.add_textbox(l, t, w, h)
        tf = tx.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(10)
            p.space_after = Pt(6)
            icon = p.add_run()
            icon.text = _BULLET_ICONS[i % len(_BULLET_ICONS)] + "  "
            icon.font.name = "Microsoft YaHei"
            icon.font.size = Pt(sz)
            icon.font.bold = True
            icon.font.color.rgb = RGBColor(*icon_color)
            body = p.add_run()
            body.text = b
            body.font.name = "Microsoft YaHei"
            body.font.size = Pt(sz)
            body.font.color.rgb = RGBColor(*_DARK)

    def _set_notes(slide, text):
        if not text:
            return
        try:
            slide.notes_slide.notes_text_frame.text = text
        except Exception:
            pass

    def _footer(slide, page_num, total):
        _text(slide, Inches(0.6), SH - Inches(0.45), Inches(4), Inches(0.3),
              "Agent-Pilot v13 · 飞书 AI 校园挑战赛", sz=9, rgb=_MUTED)
        _text(slide, SW - Inches(1.5), SH - Inches(0.45), Inches(1.2), Inches(0.3),
              f"{page_num} / {total}", sz=9, rgb=_MUTED, align="right")

    total_pages = len(outline)

    # ═══════ SLIDE 1: COVER ═══════
    cover = prs.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = RGBColor(*_NAVY)
    _rect(cover, 0, 0, Inches(0.5), SH, _GOLD)
    _rect(cover, 0, SH - Inches(0.15), SW, Inches(0.15), _GOLD)
    _rect(cover, Inches(0.5), 0, Inches(0.08), SH, _NAVY_MID)
    _text(cover, Inches(1.2), Inches(1.8), Inches(11), Inches(1.6),
          title, sz=52, bold=True, rgb=_WHITE)
    subtitle = (outline[0].get("bullets") or [""])[0] if outline else ""
    if subtitle:
        _rect(cover, Inches(1.2), Inches(3.7), Inches(3), Inches(0.04), _GOLD)
        _text(cover, Inches(1.2), Inches(3.9), Inches(11), Inches(0.7),
              subtitle, sz=22, rgb=_GOLD_LIGHT)
    _text(cover, Inches(1.2), Inches(5.8), Inches(5), Inches(0.4),
          "Agent-Pilot · 飞书 AI 校园挑战赛", sz=14, rgb=_MUTED)
    _text(cover, Inches(1.2), Inches(6.3), Inches(5), Inches(0.4),
          "戴尚好 / 李洁盈", sz=14, rgb=_GOLD_LIGHT)
    _set_notes(cover, outline[0].get("note", "") if outline else "")

    # ═══════ SLIDE 2: TOC (if outline[1] is a 目录) ═══════
    content_start = 1
    if len(outline) > 2 and "目录" in (outline[1].get("title") or ""):
        toc = prs.slides.add_slide(blank)
        toc.background.fill.solid()
        toc.background.fill.fore_color.rgb = RGBColor(*_WHITE)
        _rect(toc, 0, 0, SW, Inches(1.6), _NAVY)
        _text(toc, Inches(0.8), Inches(0.4), Inches(10), Inches(1.0),
              outline[1].get("title", "目录"), sz=36, bold=True, rgb=_WHITE)
        _rect(toc, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.04), _GOLD)
        toc_bullets = outline[1].get("bullets") or []
        cols = min(len(toc_bullets), 3)
        col_w = Inches(3.8)
        for idx, item in enumerate(toc_bullets):
            col = idx % cols
            row = idx // cols
            x = Inches(0.8) + col * col_w
            y = Inches(2.2) + row * Inches(1.4)
            _rect(toc, x, y, Inches(3.4), Inches(1.1), _LIGHT)
            _rect(toc, x, y, Inches(0.08), Inches(1.1), _SECTION_COLORS[idx % len(_SECTION_COLORS)])
            num_text = f"0{idx + 1}" if idx < 9 else str(idx + 1)
            _text(toc, x + Inches(0.3), y + Inches(0.15), Inches(0.6), Inches(0.5),
                  num_text, sz=24, bold=True, rgb=_SECTION_COLORS[idx % len(_SECTION_COLORS)])
            _text(toc, x + Inches(1.0), y + Inches(0.25), Inches(2.2), Inches(0.6),
                  item, sz=16, rgb=_DARK)
        _footer(toc, 2, total_pages)
        _set_notes(toc, outline[1].get("note", ""))
        content_start = 2

    # ═══════ CONTENT SLIDES ═══════
    content_pages = outline[content_start:-1] if len(outline) > content_start + 1 else outline[content_start:]
    for ci, page in enumerate(content_pages):
        slide_num = content_start + ci + 1
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*_WHITE)
        section_color = _SECTION_COLORS[ci % len(_SECTION_COLORS)]
        page_title = page.get("title", "")
        bullets = page.get("bullets") or []
        use_split = len(bullets) >= 3 and ci % 3 == 0
        use_cards = len(bullets) >= 3 and ci % 3 == 1

        # top nav bar
        _rect(slide, 0, 0, SW, Inches(0.06), section_color)
        # left accent strip
        _rect(slide, 0, 0, Inches(0.35), SH, _NAVY)
        _rect(slide, Inches(0.35), 0, Inches(0.06), SH, section_color)

        # section number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), Inches(0.4), Inches(0.7), Inches(0.7),
        )
        circle.line.fill.background()
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*section_color)
        _text(slide, Inches(0.7), Inches(0.48), Inches(0.7), Inches(0.6),
              f"{slide_num - 1}", sz=22, bold=True, rgb=_WHITE, align="center")

        # title
        _text(slide, Inches(1.6), Inches(0.4), Inches(10), Inches(0.8),
              page_title, sz=30, bold=True, rgb=_NAVY)
        _rect(slide, Inches(1.6), Inches(1.2), Inches(1.5), Inches(0.04), section_color)

        if use_cards and len(bullets) >= 3:
            # card layout: 3 across
            card_w = Inches(3.6)
            gap = Inches(0.3)
            start_x = Inches(0.8)
            for bi, b in enumerate(bullets[:6]):
                col = bi % 3
                row = bi // 3
                cx = start_x + col * (card_w + gap)
                cy = Inches(1.8) + row * Inches(2.4)
                _rect(slide, cx, cy, card_w, Inches(2.1), _LIGHT)
                _rect(slide, cx, cy, card_w, Inches(0.06), section_color)
                _text(slide, cx + Inches(0.3), cy + Inches(0.25), card_w - Inches(0.6), Inches(0.4),
                      _BULLET_ICONS[bi % len(_BULLET_ICONS)], sz=28, rgb=section_color, bold=True)
                _text(slide, cx + Inches(0.3), cy + Inches(0.75), card_w - Inches(0.6), Inches(1.2),
                      b, sz=15, rgb=_DARK)
        elif use_split:
            # split: left bullets, right highlight box
            _numbered_bullets(slide, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5),
                              bullets, accent=section_color, sz=18)
            _rect(slide, Inches(8.0), Inches(1.8), Inches(4.8), Inches(4.5), _LIGHT)
            _rect(slide, Inches(8.0), Inches(1.8), Inches(0.06), Inches(4.5), section_color)
            highlight = bullets[0] if bullets else page_title
            _text(slide, Inches(8.4), Inches(2.4), Inches(4.0), Inches(1.0),
                  "核心观点", sz=14, bold=True, rgb=section_color)
            _text(slide, Inches(8.4), Inches(3.0), Inches(4.0), Inches(2.5),
                  highlight, sz=16, rgb=_DARK)
        else:
            # standard: icon bullets full width
            _icon_bullets(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8),
                          bullets, icon_color=section_color, sz=18)

        _footer(slide, slide_num, total_pages)
        _set_notes(slide, page.get("note", ""))

    # ═══════ LAST SLIDE: THANK YOU ═══════
    thank = outline[-1] if len(outline) > 1 else {"title": "Thank You", "bullets": [], "note": ""}
    last = prs.slides.add_slide(blank)
    last.background.fill.solid()
    last.background.fill.fore_color.rgb = RGBColor(*_NAVY)
    _rect(last, 0, 0, Inches(0.5), SH, _GOLD)
    _rect(last, 0, SH - Inches(0.15), SW, Inches(0.15), _GOLD)

    _text(last, Inches(1.0), Inches(2.2), Inches(11), Inches(1.4),
          thank.get("title", "Thank You"), sz=56, bold=True, rgb=_WHITE, align="center")
    _rect(last, Inches(5.5), Inches(3.8), Inches(2.5), Inches(0.04), _GOLD)
    extras = thank.get("bullets") or []
    if extras:
        _text(last, Inches(1.0), Inches(4.2), Inches(11), Inches(2.0),
              "\n".join(extras), sz=16, rgb=_GOLD_LIGHT, align="center")
    _text(last, Inches(1.0), Inches(6.4), Inches(11), Inches(0.4),
          "戴尚好 / 李洁盈 · Agent-Pilot · 飞书 AI 校园挑战赛", sz=13, rgb=_MUTED, align="center")
    _set_notes(last, thank.get("note", "致谢评委"))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return len(prs.slides)


# ── Slidev export ─────────────────────────────────────────────────────────────


def _outline_to_slidev_md(title: str, outline: List[Dict[str, Any]]) -> str:
    parts = [
        "---",
        "theme: seriph",
        f"title: {title}",
        "class: text-center",
        "highlighter: shiki",
        "lineNumbers: false",
        "---",
        "",
        f"# {title}",
        "",
        "Agent-Pilot · 从 IM 对话到演示稿的一键智能闭环",
        "",
    ]
    for page in outline:
        parts.append("---")
        parts.append("")
        parts.append(f"# {page.get('title', '')}")
        parts.append("")
        for b in page.get("bullets", []):
            parts.append(f"- {b}")
        if page.get("note"):
            parts.append("")
            parts.append("<!-- ")
            parts.append(page["note"])
            parts.append("-->")
        parts.append("")
    return "\n".join(parts)


def _build_slidev_html(slidev_md: Path, out_dir: Path) -> Optional[Path]:
    """Try to build Slidev HTML if `npx` is available; else skip silently."""
    npx = shutil.which("npx")
    if not npx:
        logger.info("slidev: npx not on PATH, skipping HTML build (md still saved)")
        return None
    try:
        result = subprocess.run(
            [npx, "--yes", "@slidev/cli", "build", str(slidev_md), "--out", str(out_dir)],
            capture_output=True, text=True, timeout=180,
        )
        if result.returncode != 0:
            logger.warning("slidev build failed: %s", result.stderr[:300])
            return None
        idx = out_dir / "index.html"
        if idx.exists():
            return idx
    except Exception as e:
        logger.warning("slidev build exception: %s", e)
    return None


# ── Speaker notes ─────────────────────────────────────────────────────────────


def _outline_to_speaker_notes(title: str, outline: List[Dict[str, Any]]) -> str:
    lines = [f"# {title} · 演讲稿\n"]
    for i, page in enumerate(outline, 1):
        lines.append(f"## 第 {i} 页：{page.get('title', '')}\n")
        note = page.get("note", "").strip()
        if note:
            lines.append(note)
        else:
            lines.append("（请用自己的语言讲 40-60 秒）")
        lines.append("")
        bullets = page.get("bullets") or []
        if bullets:
            lines.append("**幻灯片要点：**")
            for b in bullets:
                lines.append(f"- {b}")
            lines.append("")
    return "\n".join(lines)


def _speaker_note_for_page(page: Dict[str, Any]) -> str:
    if page.get("note"):
        return page["note"]
    title = page.get("title", "")
    bullets = page.get("bullets", [])
    if not bullets:
        return f"本页主题：{title}"
    return f"本页主题：{title}。要点：" + "；".join(bullets[:3]) + "。"


def _make_tts_optional(title: str, outline: List[Dict[str, Any]],
                       out_dir: Path, slide_id: str) -> Optional[Path]:
    """Best-effort TTS. Disabled by default (network dependency)."""
    if os.getenv("AGENT_PILOT_ENABLE_TTS", "0") not in ("1", "true", "yes"):
        return None
    try:
        from gtts import gTTS
    except ImportError:
        logger.info("tts: gTTS not installed, skipping mp3")
        return None
    try:
        narration = "\n".join(p.get("note") or p.get("title", "") for p in outline)
        if not narration.strip():
            return None
        mp3_path = out_dir / f"{slide_id}.mp3"
        gTTS(text=narration[:4500], lang="zh-cn").save(str(mp3_path))
        return mp3_path
    except Exception as e:
        logger.warning("tts mp3 generation failed: %s", e)
        return None


def _parse_slidev_md(path: Path) -> List[Dict[str, Any]]:
    try:
        raw = path.read_text(encoding="utf-8")
    except Exception:
        return []
    pages: List[Dict[str, Any]] = []
    for chunk in raw.split("\n---\n"):
        lines = [l for l in chunk.splitlines() if l.strip()]
        if not lines:
            continue
        title = ""
        bullets: List[str] = []
        for ln in lines:
            if ln.startswith("# "):
                title = ln[2:].strip()
            elif ln.startswith("- "):
                bullets.append(ln[2:].strip())
        if title:
            pages.append({"title": title, "bullets": bullets, "note": ""})
    return pages


# ── Feishu integration ───────────────────────────────────────────────────────


def _create_feishu_preview_doc(title: str, outline: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Create a Feishu Docx that mirrors the deck for mobile-first viewing."""
    try:
        from config import Config
        if not (Config.FEISHU_APP_ID and Config.FEISHU_APP_SECRET):
            return {}
        import lark_oapi.api.docx.v1 as docx_api
        from bot.feishu_client import get_client

        client = get_client()
        doc_title = f"[演示稿预览] {title}"
        req = (
            docx_api.CreateDocumentRequest.builder()
            .request_body(docx_api.CreateDocumentRequestBody.builder().title(doc_title).build())
            .build()
        )
        resp = client.docx.v1.document.create(req)
        if not resp.success() or not resp.data or not resp.data.document:
            logger.warning("preview doc create failed: code=%s msg=%s",
                           getattr(resp, "code", "?"), getattr(resp, "msg", "?"))
            return {}
        doc_token = resp.data.document.document_id
        domain = getattr(Config, "FEISHU_TENANT_DOMAIN", "") or "rcnqvnspd31b.feishu.cn"

        blocks = _outline_to_docx_blocks(outline)
        if blocks:
            from lark_oapi.api.docx.v1 import (
                CreateDocumentBlockChildrenRequest,
                CreateDocumentBlockChildrenRequestBody,
            )
            BATCH = 50
            for i in range(0, len(blocks), BATCH):
                req2 = (
                    CreateDocumentBlockChildrenRequest.builder()
                    .document_id(doc_token)
                    .block_id(doc_token)
                    .request_body(
                        CreateDocumentBlockChildrenRequestBody.builder()
                        .children(blocks[i:i + BATCH]).build()
                    ).build()
                )
                client.docx.v1.document_block_children.create(req2)
        return {"doc_token": doc_token, "url": f"https://{domain}/docx/{doc_token}"}
    except Exception as e:
        logger.warning("preview doc creation failed: %s", e)
        return {}


def _outline_to_docx_blocks(outline: List[Dict[str, Any]]) -> list:
    try:
        from lark_oapi.api.docx.v1 import Block, Text, TextElement, TextRun
    except ImportError:
        return []

    def _tb(text: str, bt: int):
        run = TextRun.builder().content(text).build()
        el = TextElement.builder().text_run(run).build()
        txt = Text.builder().elements([el]).build()
        try:
            bb = Block.builder().block_type(bt)
            if bt == 3:
                return bb.heading1(txt).build()
            if bt == 4:
                return bb.heading2(txt).build()
            if bt == 5:
                return bb.heading3(txt).build()
            if bt == 12:
                return bb.bullet(txt).build()
            return bb.text(txt).build()
        except Exception:
            return Block.builder().block_type(2).text(txt).build()

    blocks = []
    for i, page in enumerate(outline, 1):
        blocks.append(_tb(f"第 {i} 页 · {page.get('title', '')}", 3))
        for b in page.get("bullets", []) or []:
            blocks.append(_tb(b, 12))
        if page.get("note"):
            blocks.append(_tb(f"演讲稿：{page['note']}", 2))
        blocks.append(_tb("", 2))
    return blocks


def _upload_pptx_to_feishu_drive(pptx_path: Path) -> Optional[str]:
    """Upload the .pptx to Feishu Drive root folder. Returns shareable URL or None."""
    try:
        from config import Config
        if not (Config.FEISHU_APP_ID and Config.FEISHU_APP_SECRET):
            return None
        import lark_oapi as lark  # noqa: F401
        from bot.feishu_client import get_client
        from lark_oapi.api.drive.v1 import (
            UploadAllFileRequest,
            UploadAllFileRequestBody,
        )

        folder_token = os.getenv("FEISHU_FOLDER_TOKEN", "")
        client = get_client()
        file_size = pptx_path.stat().st_size  # 用 stat() 而非 read-all 获取真实大小
        fh = open(pptx_path, "rb")
        try:
            body = (
                UploadAllFileRequestBody.builder()
                .file_name(pptx_path.name)
                .parent_type("explorer")
                .parent_node(folder_token or "")
                .size(file_size)
                .file(fh)             # 传文件句柄（io stream），不是 bytes
                .build()
            )
            req = UploadAllFileRequest.builder().request_body(body).build()
            resp = client.drive.v1.file.upload_all(req)
        finally:
            fh.close()

        if not resp.success() or not resp.data:
            logger.warning("pptx upload failed: code=%s msg=%s",
                           getattr(resp, "code", "?"), getattr(resp, "msg", "?"))
            return None
        file_token = getattr(resp.data, "file_token", "") or ""
        if not file_token:
            return None
        domain = getattr(Config, "FEISHU_TENANT_DOMAIN", "") or "rcnqvnspd31b.feishu.cn"
        return f"https://{domain}/file/{file_token}"
    except Exception as e:
        logger.warning("pptx upload to drive failed: %s", e)
        return None
